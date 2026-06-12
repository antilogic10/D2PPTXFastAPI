from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os


def extract_images_from_ppt(
    ppt_path: str,
    output_folder: str = "extracted_images",
):
    prs = Presentation(ppt_path)
    slide = prs.slides[0]

    os.makedirs(output_folder, exist_ok=True)

    saved_files = []

    for idx, shape in enumerate(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue

        image = shape.image
        image_bytes = image.blob
        image_ext = image.ext  # 'png', 'jpeg', etc.

        file_name = f"{shape.name or f'image_{idx}'}.{image_ext}"
        file_path = os.path.join(output_folder, file_name)

        with open(file_path, "wb") as f:
            f.write(image_bytes)

        saved_files.append(file_path)

    return saved_files


files = extract_images_from_ppt(
    ppt_path="./p1.pptx",
    output_folder="extracted_icons"
)

for f in files:
    print("Saved:", f)
