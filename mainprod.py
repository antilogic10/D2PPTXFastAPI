from fastapi import FastAPI, HTTPException, File, UploadFile
from fastapi.responses import HTMLResponse
from pydantic import BaseModel
import os
from pptx import Presentation
from google import genai
import json
import requests
import tempfile
from datetime import datetime
from fastapi.middleware.cors import CORSMiddleware
import uuid
from fastapi.staticfiles import StaticFiles
import shutil
from typing import List
import copy

UPLOAD_DIR = "uploaded_files"
GENERATED_DIR = "generated_files"
DOMAIN_NAME = os.getenv("DOMAIN_NAME", "http://localhost:8000")
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.umask(0o022)

app = FastAPI(root_path="/api")

# Allow all origins
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # <-- allows all origins
    allow_credentials=True,
    allow_methods=["*"],  # <-- allows all methods (GET, POST, etc.)
    allow_headers=["*"],  # <-- allows all headers
)

app.mount(f"/{GENERATED_DIR}", StaticFiles(directory=GENERATED_DIR), name="files")
# Serve uploaded files
app.mount(f"/{UPLOAD_DIR}", StaticFiles(directory=UPLOAD_DIR), name="files")
# Pydantic model for request body
class PPTRequest(BaseModel):
    fileUrl: str   # Name of the pptx template file
    content: str  # Unstructured content to be filled in the pptx
    imageUrl: str   # image url uploaded to gemini for context
    rewriteWithAi: bool = False  # Whether to rewrite content with AI
# Initialize Gemini client
client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))

def list_text_boxes(pptx_path: str, slide_index: int):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    placeholders = {}

    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            # Check if any paragraph is bulleted
            is_list = any(p.level > 0 or p.text.strip().startswith("•") for p in shape.text_frame.paragraphs)

            # Use the first non-empty run text as the "placeholder key"
            placeholder_key = shape.text.strip()

            if is_list:
                items = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                placeholders[placeholder_key] = {"type": "list", "items": items}
            else:
                placeholders[placeholder_key] = {"type": "text", "value": shape.text.strip()}

    return placeholders


def updateTemplatePlaceholders(pptx_path: str, slide_index: int, replacements: dict):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]

    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            original_text = shape.text.strip()
            if original_text in replacements:
                if isinstance(replacements[original_text], dict) and "value" in replacements[original_text]:
                    new_value = replacements[original_text]["value"]
                else:
                    new_value = replacements[original_text]
                
                # Detect if template shape is a list
                is_list_shape = (
                    len(shape.text_frame.paragraphs) > 1
                    or any(p.level > 0 for p in shape.text_frame.paragraphs)
                    or any(p._pPr is not None and p._pPr.xpath(".//a:buChar") for p in shape.text_frame.paragraphs)
                    or any(p._pPr is not None and p._pPr.xpath(".//a:buAutoNum") for p in shape.text_frame.paragraphs)
                )

                # 🔹 If template expects a list but Gemini returned string → wrap in list
                if is_list_shape and isinstance(new_value, str):
                    new_value = [new_value]

                # 🔹 If template expects plain text but Gemini returned list → join
                if not is_list_shape and isinstance(new_value, list):
                    new_value = " ".join(new_value)

                # --- Replace text based on detected type ---
                if isinstance(new_value, str):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip() == original_text:
                                run.text = new_value
                
                elif isinstance(new_value, type(None)):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip() == original_text:
                                run.text = ""

                elif isinstance(new_value, list):
                    paragraphs = shape.text_frame.paragraphs
                    template_p = paragraphs[0] if len(paragraphs) > 0 else None

                    counter = 0
                    for item in new_value:
                        if counter < len(paragraphs):
                            p = paragraphs[counter]
                            if p.runs:
                                p.runs[0].text = item
                                for r in p.runs[1:]:
                                    r.text = ""
                            else:
                                p.text = item
                        else:
                            p = shape.text_frame.add_paragraph()
                            p.text = item

                            if template_p is not None:
                                # ---- Copy paragraph-level formatting (bullets, indent, alignment)
                                p.level = template_p.level
                                if template_p._pPr is not None:
                                    try:
                                        pPr_copy = copy.deepcopy(template_p._pPr)
                                        # remove old paragraph properties if they exist
                                        if p._pPr is not None:
                                            p._p.remove(p._pPr)
                                        # inject cloned paragraph properties
                                        p._p.insert(0, pPr_copy)
                                    except Exception as e:
                                        print("Error copying paragraph properties:", e)

                                # ---- Copy text run formatting (font, bold, italic, color, size)
                                if template_p.runs:
                                    t_run = template_p.runs[0]
                                    if p.runs:
                                        run = p.runs[0]
                                    else:
                                        run = p.add_run()

                                    run.font.bold = t_run.font.bold
                                    run.font.italic = t_run.font.italic
                                    run.font.size = t_run.font.size
                                    run.font.name = t_run.font.name
                        counter += 1

                    # ---- Clear extra bullets / paragraphs
                    for p in paragraphs[counter:]:
                        if p.runs:
                            for r in p.runs:
                                r.text = ""
                        else:
                            p.text = ""

                else:
                    print(f"Skipping unknown type for {original_text}")

    output_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx").name
    prs.save(output_path)
    return output_path


def download_pptx(url: str) -> str:
    # """Download PPTX from the given URL and save locally"""
    response = requests.get(url)
    if response.status_code != 200:
        raise HTTPException(status_code=400, detail="Could not download PPT file")
    
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp_file.write(response.content)
    tmp_file.close()
    os.chmod(tmp_file.name, 0o755)
    return tmp_file.name

def download_image(url: str) -> str:
    # """Download image from the given URL and save locally"""
    response = requests.get(url)
    ext=url.split('.')[-1] if '.' in url else ''

    if response.status_code != 200:
        raise HTTPException(status_code=400, detail="Could not download image file")
    
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}")
    tmp_file.write(response.content)
    tmp_file.close()
    os.chmod(tmp_file.name, 0o755)
    return tmp_file.name

STARTED_AT = datetime.utcnow()

@app.get("/",response_class=HTMLResponse)
def home():
    uptime = datetime.utcnow() - STARTED_AT
    
    # 🧮 Get Disk Usage
    total, used, free = shutil.disk_usage("/")
    used_gb = used / (2**30)
    total_gb = total / (2**30)
    free_gb = free / (2**30)
    percent_used = (used / total) * 100

    # 🧱 Return HTML
    return f"""
    <!doctype html>
    <html lang="en">
    <head>
        <meta charset="utf-8"/>
        <title>{app.title} • Status</title>
        <style>
            body {{
                font-family: system-ui, sans-serif;
                background: #0f172a;
                color: #f1f5f9;
                display: flex;
                justify-content: center;
                align-items: center;
                height: 100vh;
                margin: 0;
            }}
            .card {{
                background: #1e293b;
                padding: 2rem 3rem;
                border-radius: 1rem;
                text-align: center;
                box-shadow: 0 10px 20px rgba(0,0,0,0.5);
                width: 400px;
            }}
            h1 {{
                margin: 0 0 0.5rem;
                font-size: 1.8rem;
                color: #38bdf8;
            }}
            p {{ margin: 0.5rem 0; color: #cbd5e1; }}
            .uptime {{
                font-size: 0.9rem;
                color: #94a3b8;
            }}
            a {{
                color: #38bdf8;
                text-decoration: none;
            }}
            a:hover {{ text-decoration: underline; }}
            .progress {{
                background: #334155;
                border-radius: 10px;
                overflow: hidden;
                height: 16px;
                width: 100%;
                margin: 10px 0;
                box-shadow: inset 0 1px 3px rgba(0,0,0,0.3);
            }}
            .progress-bar {{
                height: 100%;
                background: linear-gradient(90deg, #38bdf8, #0ea5e9);
                width: {percent_used:.2f}%;
                transition: width 0.5s ease-in-out;
            }}
            .disk-info {{
                font-size: 0.85rem;
                color: #a1a1aa;
            }}
        </style>
    </head>
    <body>
        <div class="card">
            <h1>🚀 {app.title} is Live</h1>
            <p>All APIs are up and working correctly.</p>
            <p class="uptime">Uptime: {uptime}</p>

            <div style="margin-top:1rem;">
                <p>💾 Disk Usage</p>
                <div class="progress">
                    <div class="progress-bar"></div>
                </div>
                <p class="disk-info">
                    Used: {used_gb:.2f} GB / {total_gb:.2f} GB<br/>
                    Free: {free_gb:.2f} GB ({100 - percent_used:.2f}%)
                </p>
            </div>

            <p style="margin-top:1rem;">
                <a href="/docs">Interactive Docs</a> • 
                <a href="/redoc">ReDoc</a>
            </p>
        </div>
    </body>
    </html>
    """

def validateJson(cleaned_json, textBoxList):
    # Check for explicit error
    if "error" in cleaned_json:
        print("\n-----\nError found in JSON:", cleaned_json["error"], "\n-----\n")
        return False

    # Check placeholder mismatch
    if len(cleaned_json.keys()) != len(textBoxList):
        print("\n-----\nPlaceholder count mismatch:", len(cleaned_json.keys()), "vs", len(textBoxList), "\n-----\n")
        return False

    # Validate values
    seen_values = set()
    for k, v in cleaned_json.items():
        if not v:
            print("\n-----\nEmpty value for key:", k, "\n-----\n")
            return False

        # Skip numeric placeholders (like 01, 02…)
        if k.isdigit():
            continue

        # Disallow "Heading 3": "Heading 3" or "Slide title": "Slide title"
        if isinstance(v, str) and v.strip().lower() == k.strip().lower():
            print("\n-----\nRepeated value for key:", k, "\n-----\n")
            return False

        # Handle unhashable types safely
        if isinstance(v, (dict, list)):
            v_hash = json.dumps(v, sort_keys=True)  # Convert to string for hashing
        else:
            v_hash = str(v).strip()

        # Optional: disallow duplicate non-numeric values
        if v_hash in seen_values:
            print("\n-----\nDuplicate value detected:", v, "\n-----\n")
            return False
        seen_values.add(v_hash)

    return True

@app.post("/generate-ppt")
def generate_ppt(req: PPTRequest):
    # Step 1: Download template
    pptx_path = download_pptx(req.fileUrl)
    textBoxList = list_text_boxes(pptx_path, 0)

    prompt = f""" """
    
    if(req.rewriteWithAi):
        prompt = f"""
        You are an expert PowerPoint slide content writer and layout-aware editor. 
        Your task is to enhance and professionally rewrite the given content so it fits clearly and neatly into the provided PowerPoint placeholders, 
        keeping the slide visually balanced and non-repetitive.
        
        If you cannot produce a valid mapping for every placeholder,
        return only this JSON:
        {{"error": "Content too short for the template. Please provide more detailed content."}}
        
        ### Inputs
        - Content: {req.content}
        - Placeholders: {json.dumps(textBoxList, indent=2)}
        
        ### Objectives
        1. Rewrite the provided content in a concise, business-professional tone.
        2. You may **add small, relevant details or context** only if it helps clarify or complete ideas, but never invent unrelated or misleading information.
        3. Do not copy identical text across multiple placeholders unless it is genuinely meant to repeat (e.g., a shared title).
        4. Keep wording compact enough so text fits inside each placeholder box — imagine a standard PowerPoint layout where 4–6 bullet points per box is ideal.
        5. The template image is **only for reference** to understand approximate space and structure. Do not infer color, shape, or visual design from it.
        
        ### Mapping Logic
        1. Determine the purpose of each placeholder (e.g., title, subtitle, step, description, list).
        2. Split and map the rewritten content logically:
           - Assign distinct yet contextually linked text to each placeholder.
           - For sequential steps (e.g., “Discover”, “Plan”, “Create”, “Deliver”), ensure each step has its own unique focus and description.
        3. If placeholder type is "list", return an array of bullet points (minimum 1 and maximum as per the template image provided). When creating bullet lists, use clean text without adding any extra symbols such as hyphens (-), asterisks (*), or other bullet markers — return plain text items only. If placeholder type is "text", return a single concise string.

        4. If any placeholder cannot be filled meaningfully, stop and return:
           {{"error": "Content too short for the template. Please provide more detailed content."}}
        
        ### Output Requirements
        - Return **only a valid JSON object**.
        - Keys = exact placeholder text from the provided list.
        - Values = strings or string arrays depending on placeholder type.
        - No explanations, markdown, or extra commentary.
        - Ensure the JSON is syntactically valid.
        """
    else:
        prompt = f"""You are a highly precise PowerPoint content mapper. 
        Your task is to map the provided content directly to the given placeholders **exactly as written**, without rewriting or rephrasing it — 
        but with logical splitting and proper assignment based on context.
        
        If you cannot produce a valid mapping for every placeholder, 
        return only this JSON:
        {{"error": "Content too short for the template. Please provide more detailed content."}}
        
        ### Inputs
        - Content: {req.content}
        - Placeholders: {json.dumps(textBoxList, indent=2)}
        
        ### Mapping Rules
        1. **Do not rewrite or rephrase** the text; only split or assign it logically.
        2. Identify the intent of each placeholder (e.g., “Discover”, “Plan”, “Create”, “Deliver”).
        3. Group related sentences, phrases, or bullet points from the content and assign them to the most contextually relevant placeholder.
           - Example: All lines mentioning “research”, “analysis”, or “identifying needs” → map to “Discover”.
           - Lines about “strategy”, “planning”, “goal setting” → map to “Plan”.
           - Lines about “design”, “development”, “execution” → map to “Create”.
           - Lines about “testing”, “delivery”, “measurement”, “results” → map to “Deliver”.
        4. If placeholder type is "list", return an array of bullet points (minimum 1 and maximum as per the template image provided). When creating bullet lists, use clean text without adding any extra symbols such as hyphens (-), asterisks (*), or other bullet markers — return plain text items only. If placeholder type is "text", return a single concise string.
        5. **Never duplicate the same sentences** across placeholders unless the content itself repeats them exactly.
        6. Maintain the factual meaning and original order of ideas wherever possible.
        7. If any placeholder cannot be filled with meaningful data, stop and return:
           {{"error": "Content too short for the template. Please provide more detailed content."}}
        8. If the provided content clearly and meaningfully fills only a subset of placeholders (for example, a title and four main sections), this is acceptable. Do not force-fill empty placeholders with guesses or duplicated text. Only leave placeholders empty if no relevant content exists for them.
        ### Output Format
        - Output **strictly valid JSON only**.
        - Keys = exact placeholder text from the provided list.
        - Values = strings or string arrays depending on placeholder type.
        - No markdown, no explanations, no extra commentary.
        """

    uploadedFile = client.files.upload(file=download_image(req.imageUrl))
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=[prompt,uploadedFile]
    )
    cleanedJson = json.loads((response.text.strip("`")).replace("json","",1).strip())
    print("\n----- Prompted ----",prompt,"\n---end prompt---","\n------\nGenerated JSON:", cleanedJson,"\n------\n")
    if not validateJson(cleanedJson, textBoxList):
        if os.path.exists(pptx_path):
            os.remove(pptx_path)
        return {"error": "Error Generating PPTX, Content too short for the template. Please provide more detailed content."}
    else:
        updated_pptx =updateTemplatePlaceholders(pptx_path, 0, cleanedJson)

        # Step 3: Generate unique filename
        unique_id = uuid.uuid4().hex[:8]  # short UUID
        public_filename = f"presentation_{unique_id}.pptx"
        public_path = os.path.join(GENERATED_DIR, public_filename)

        # Move to public folder
        shutil.copy(updated_pptx, public_path)
        # make the file publicly readable
        os.chmod(public_path, 0o755)
        # Delete the temporary file
        if os.path.exists(updated_pptx):
            os.remove(updated_pptx)

        if os.path.exists(pptx_path):
            os.remove(pptx_path)

        # Step 4: Return public URL
        file_url = f"{DOMAIN_NAME}api/{GENERATED_DIR}/{public_filename}"
        return {"file_url": file_url}

@app.post("/upload-files/")
async def upload_files(files: List[UploadFile] = File(...)):
    saved_files = []

    for file in files:
        # Always use the original filename
        filename = file.filename
        file_path = os.path.join(UPLOAD_DIR, filename)

        # "wb" mode automatically replaces file if it already exists
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Make the file publicly readable
        os.chmod(file_path, 0o755)

        # Build file URL
        file_url = f"{DOMAIN_NAME}api/{UPLOAD_DIR}/{filename}"
        saved_files.append({"filename": filename, "url": file_url})

    return {"uploaded": saved_files}
