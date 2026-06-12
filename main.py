from fastapi import FastAPI, HTTPException, File, UploadFile
from fastapi.responses import HTMLResponse
from pydantic import BaseModel
import os
from pptx import Presentation
import json
import requests
import tempfile
from datetime import datetime
from fastapi.middleware.cors import CORSMiddleware
import uuid
from fastapi.staticfiles import StaticFiles
import shutil
from typing import List
from pinecone import Pinecone
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
from copy import deepcopy
from io import BytesIO
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()

UPLOAD_DIR = "uploaded_files"
GENERATED_DIR = "generated_files"
DOMAIN_NAME = os.getenv("DOMAIN_NAME", "http://localhost:8000")
INDEX_NAME = "icons-store"
ICONS_NAMESPACE = "icons-namespace"
BATCH_SIZE = 50
BASE_DIR = Path(__file__).resolve().parent
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.umask(0o022)

app = FastAPI(root_path="/api")
pc = Pinecone(api_key=os.getenv("PINECONE_KEY", "pcsk_5dAecR_ThUqqVrC9cYhb3YbJXjVz8sDp1zjraPPhrGQ1j1Y3L3BzCrPmgCFLKPbmnjpXnR"))
dense_index = pc.Index(INDEX_NAME)

# Allow all origins
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # <-- allows all origins
    allow_credentials=True,
    allow_methods=["*"],  # <-- allows all methods (GET, POST, etc.)
    allow_headers=["*"],  # <-- allows all headers
)

app.mount(f"/{GENERATED_DIR}", StaticFiles(directory=GENERATED_DIR), name="files")
# Serve uploaded files
app.mount(f"/{UPLOAD_DIR}", StaticFiles(directory=UPLOAD_DIR), name="files")
# Pydantic model for request body
class PPTRequest(BaseModel):
    fileUrl: str   # Name of the pptx template file
    content: str  # Unstructured content to be filled in the pptx
    imageUrl: str   # image url uploaded to gemini for context
    rewriteWithAi: bool = False  # Whether to rewrite content with AI

class ConsolidateSlidesRequest(BaseModel):
    ppt_paths: List[str]

# Initialize OPENAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def copy_slide(source_slide, target_prs):
    blank_layout = target_prs.slide_layouts[6]
    target_slide = target_prs.slides.add_slide(blank_layout)

    for shape in source_slide.shapes:

        # ---- IMAGES ----
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_stream = BytesIO(shape.image.blob)

            target_slide.shapes.add_picture(
                image_stream,
                shape.left,
                shape.top,
                shape.width,
                shape.height
            )

        # ---- TEXT / AUTOSHAPES / EVERYTHING ELSE ----
        else:
            new_el = deepcopy(shape.element)
            target_slide.shapes._spTree.insert_element_before(
                new_el, 'p:extLst'
            )


def list_text_boxes(pptx_path, slide_index: int):
    prs = pptx_path #Presentation(pptx_path)
    slide = prs.slides[slide_index]
    placeholders = {}

    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            # Check if any paragraph is bulleted
            is_list = any(p.level > 0 or p.text.strip().startswith("•") for p in shape.text_frame.paragraphs)

            # Use the first non-empty run text as the "placeholder key"
            placeholder_key = shape.text.strip()

            if is_list:
                items = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                placeholders[placeholder_key] = {"type": "list", "items": items}
            else:
                placeholders[placeholder_key] = {"type": "text", "value": shape.text.strip()}

    return placeholders


def updateTemplatePlaceholders(pptx_path, slide_index: int, replacements: dict):
    prs = pptx_path #Presentation(pptx_path)
    slide = prs.slides[slide_index]

    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            original_text = shape.text.strip()
            if original_text in replacements:
                if isinstance(replacements[original_text], dict) and "value" in replacements[original_text]:
                    new_value = replacements[original_text]["value"]
                else:
                    new_value = replacements[original_text]
                
                # Detect if template shape is a list
                is_list_shape = (
                    len(shape.text_frame.paragraphs) > 1
                    or any(p.level > 0 for p in shape.text_frame.paragraphs)
                    or any(p._pPr is not None and p._pPr.xpath(".//a:buChar") for p in shape.text_frame.paragraphs)
                    or any(p._pPr is not None and p._pPr.xpath(".//a:buAutoNum") for p in shape.text_frame.paragraphs)
                )

                # 🔹 If template expects a list but Gemini returned string → wrap in list
                if is_list_shape and isinstance(new_value, str):
                    new_value = [new_value]

                # 🔹 If template expects plain text but Gemini returned list → join
                if not is_list_shape and isinstance(new_value, list):
                    new_value = " ".join(new_value)

                # --- Replace text based on detected type ---
                if isinstance(new_value, str):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip() == original_text:
                                run.text = new_value
                
                elif isinstance(new_value, type(None)):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip() == original_text:
                                run.text = ""

                elif isinstance(new_value, list):
                    paragraphs = shape.text_frame.paragraphs

                    counter = 0
                    for item in new_value:
                        if counter < len(paragraphs):
                            # ✅ Replace only the text of the first run, preserve formatting
                            if paragraphs[counter].runs:
                                paragraphs[counter].runs[0].text = item
                                # Clear out extra runs if any
                                for r in paragraphs[counter].runs[1:]:
                                    r.text = ""
                            else:
                                paragraphs[counter].text = item
                        else:
                            # ✅ If template doesn't have enough list items, add new ones
                            p = shape.text_frame.add_paragraph()
                            p.text = item
                            p.level = 0
                        counter += 1

                    # ✅ Clear any extra template bullets beyond what Gemini gave
                    for p in paragraphs[counter:]:
                        if p.runs:
                            for r in p.runs:
                                r.text = ""
                        else:
                            p.text = ""

                else:
                    print(f"Skipping unknown type for {original_text}")

    output_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx").name
    prs.save(output_path)
    return output_path


def download_pptx(url: str) -> str:
    # """Download PPTX from the given URL and save locally"""
    response = requests.get(url)
    if response.status_code != 200:
        raise HTTPException(status_code=400, detail="Could not download PPT file")
    
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp_file.write(response.content)
    tmp_file.close()
    os.chmod(tmp_file.name, 0o755)
    return tmp_file.name

def download_image(url: str) -> str:
    """
    Downloads an image and converts it into a base64 data URL
    usable directly with OpenAI Vision APIs.
    """

    import requests
    import base64
    from fastapi import HTTPException

    response = requests.get(url)

    if response.status_code != 200:
        raise HTTPException(
            status_code=400,
            detail="Could not download image file"
        )

    # Detect content type
    content_type = response.headers.get("Content-Type", "image/png")

    # Convert image bytes to base64
    base64_image = base64.b64encode(response.content).decode("utf-8")

    # Build OpenAI-compatible data URL
    return f"data:{content_type};base64,{base64_image}"

def get_icon(query: str):
    return dense_index.search(
        namespace=ICONS_NAMESPACE,
        query={
            "top_k": 2,
            "inputs": {
                'text': query
            }
        }
    )

def icons_identifier(ppt_path, slide_index: int):
    prs = ppt_path #Presentation(ppt_path)
    slide = prs.slides[slide_index]

    saved_files = []

    for idx, shape in enumerate(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        # if shape.width == shape.height:
        #     continue
        image = shape.image
        image_ext = image.ext  # 'png', 'jpeg', etc.


        file_name = f"{shape.name or f'image_{idx}'}.{image_ext}"
        saved_files.append({
            "file_name": file_name, 
            "shape": shape
            })


    return saved_files

def replace_icons(slide, slide_index, old_shape, new_image_path):
    """
    Replaces a picture shape with a new image,
    preserving position, size, and rotation.
    """
    prs = slide #Presentation(slide)
    slide = prs.slides[slide_index]

    # Capture geometry
    left = old_shape.left
    top = old_shape.top
    width = old_shape.width
    height = old_shape.height
    rotation = old_shape.rotation
    print("replacing Icon", old_shape.name,)
    print("item Icon", new_image_path,)
    # Remove old shape
    slide.shapes._spTree.remove(old_shape._element)
    imgreplacer = BASE_DIR / "png_icons" / new_image_path
    print("item Icon", imgreplacer,)
    # Add new image
    new_shape = slide.shapes.add_picture(
        str(imgreplacer),
        left,
        top,
        width=width,
        height=height
    )

    # Restore rotation
    new_shape.rotation = rotation

    return new_shape


STARTED_AT = datetime.utcnow()

@app.get("/",response_class=HTMLResponse)
def home():
    uptime = datetime.utcnow() - STARTED_AT
    
    # 🧮 Get Disk Usage
    total, used, free = shutil.disk_usage("/")
    used_gb = used / (2**30)
    total_gb = total / (2**30)
    free_gb = free / (2**30)
    percent_used = (used / total) * 100

    # 🧱 Return HTML
    return f"""
    <!doctype html>
    <html lang="en">
    <head>
        <meta charset="utf-8"/>
        <title>{app.title} • Status</title>
        <style>
            body {{
                font-family: system-ui, sans-serif;
                background: #0f172a;
                color: #f1f5f9;
                display: flex;
                justify-content: center;
                align-items: center;
                height: 100vh;
                margin: 0;
            }}
            .card {{
                background: #1e293b;
                padding: 2rem 3rem;
                border-radius: 1rem;
                text-align: center;
                box-shadow: 0 10px 20px rgba(0,0,0,0.5);
                width: 400px;
            }}
            h1 {{
                margin: 0 0 0.5rem;
                font-size: 1.8rem;
                color: #38bdf8;
            }}
            p {{ margin: 0.5rem 0; color: #cbd5e1; }}
            .uptime {{
                font-size: 0.9rem;
                color: #94a3b8;
            }}
            a {{
                color: #38bdf8;
                text-decoration: none;
            }}
            a:hover {{ text-decoration: underline; }}
            .progress {{
                background: #334155;
                border-radius: 10px;
                overflow: hidden;
                height: 16px;
                width: 100%;
                margin: 10px 0;
                box-shadow: inset 0 1px 3px rgba(0,0,0,0.3);
            }}
            .progress-bar {{
                height: 100%;
                background: linear-gradient(90deg, #38bdf8, #0ea5e9);
                width: {percent_used:.2f}%;
                transition: width 0.5s ease-in-out;
            }}
            .disk-info {{
                font-size: 0.85rem;
                color: #a1a1aa;
            }}
        </style>
    </head>
    <body>
        <div class="card">
            <h1>🚀 {app.title} is Live</h1>
            <p>All APIs are up and working correctly.</p>
            <p class="uptime">Uptime: {uptime}</p>

            <div style="margin-top:1rem;">
                <p>💾 Disk Usage</p>
                <div class="progress">
                    <div class="progress-bar"></div>
                </div>
                <p class="disk-info">
                    Used: {used_gb:.2f} GB / {total_gb:.2f} GB<br/>
                    Free: {free_gb:.2f} GB ({100 - percent_used:.2f}%)
                </p>
            </div>

            <p style="margin-top:1rem;">
                <a href="/docs">Interactive Docs</a> • 
                <a href="/redoc">ReDoc</a>
            </p>
        </div>
    </body>
    </html>
    """

def validateJson(cleaned_json, textBoxList):
    # Check for explicit error
    if "error" in cleaned_json:
        print("\n-----\nError found in JSON:", cleaned_json["error"], "\n-----\n")
        return False

    # Check placeholder mismatch
    if len(cleaned_json.keys()) != len(textBoxList):
        print("\n-----\nPlaceholder count mismatch:", len(cleaned_json.keys()), "vs", len(textBoxList), "\n-----\n")
        return False

    # Validate values
    seen_values = set()
    for k, v in cleaned_json.items():
        if not v:
            print("\n-----\nEmpty value for key:", k, "\n-----\n")
            return False

        # Skip numeric placeholders (like 01, 02…)
        if k.isdigit():
            continue

        # Disallow "Heading 3": "Heading 3" or "Slide title": "Slide title"
        if isinstance(v, str) and v.strip().lower() == k.strip().lower():
            print("\n-----\nRepeated value for key:", k, "\n-----\n")
            return False

        # Handle unhashable types safely
        if isinstance(v, (dict, list)):
            v_hash = json.dumps(v, sort_keys=True)  # Convert to string for hashing
        else:
            v_hash = str(v).strip()

        # Optional: disallow duplicate non-numeric values
        if v_hash in seen_values:
            print("\n-----\nDuplicate value detected:", v, "\n-----\n")
            return False
        seen_values.add(v_hash)

    return True

@app.post("/generate-ppt")
def generate_ppt(req: PPTRequest):
    # Step 1: Download template
    pptx_path = download_pptx(req.fileUrl)
    pptx_reference = Presentation(pptx_path)
    textBoxList = list_text_boxes(pptx_reference, 0)
    iconsCount = icons_identifier(pptx_reference, 0)

    prompt = f""" """
    
    if(req.rewriteWithAi):
        prompt = f"""
        You are an expert PowerPoint slide content writer and layout-aware editor. 
        Your task is to enhance and professionally rewrite the given content so it fits clearly and neatly into the provided PowerPoint placeholders, 
        keeping the slide visually balanced and non-repetitive. Don't over flow the given content by more that 5-10 words,
        keep the word count almost similar to the existing one.
        
        If you cannot produce a valid mapping for every placeholder,
        return only this JSON:
        {{"error": "Content too short for the template. Please provide more detailed content."}}
        
        ### Inputs
        - Content: {req.content}
        - Placeholders: {json.dumps(textBoxList, indent=2)}
        
        ### Objectives
        1. Rewrite the provided content in a concise, business-professional tone.
        2. You may **add small, relevant details or context** only if it helps clarify or complete ideas, but never invent unrelated or misleading information.
        3. Do not copy identical text across multiple placeholders unless it is genuinely meant to repeat (e.g., a shared title).
        4. Keep wording compact enough so text fits inside each placeholder box — imagine a standard PowerPoint layout where 4–6 bullet points per box is ideal.
        5. The template image is **only for reference** to understand approximate space and structure. Do not infer color, shape, or visual design from it.
        6. if iconCount->({len(iconsCount)}) is greater than 0 then for each section you generate content for also generate icon intents for each icon on the slide as per the following:
        7. There are {len(iconsCount)} icons on the slide; do not reference or describe them in the text. in the json at the end just give a key called iconIntents with a list of {len(iconsCount)} strings describing the intended meaning or concept for each icon in order for each section determining it by the heading and content you generated.
        8. if iconCount->({len(iconsCount)}) is 0 then do not add iconIntents key in the final json.
        9. if iconCount->({len(iconsCount)}) is greater than 0 then give an array of {len(iconsCount)} strings for iconIntents in a json key <iconIntents> and the rest of the content in json key <replacementContent>.
        10. if iconCount->({len(iconsCount)}) is 0 then give only the content in json key <replacementContent>.

        ### Mapping Logic
        1. Determine the purpose of each placeholder (e.g., title, subtitle, step, description, list).
        2. Split and map the rewritten content logically:
           - Assign distinct yet contextually linked text to each placeholder.
           - For sequential steps (e.g., “Discover”, “Plan”, “Create”, “Deliver”), ensure each step has its own unique focus and description.
        3. If placeholder type is "list", return an array of bullet points (minimum 1 and maximum as per the template image provided). When creating bullet lists, use clean text without adding any extra symbols such as hyphens (-), asterisks (*), or other bullet markers — return plain text items only. If placeholder type is "text", return a single concise string.

        4. If any placeholder cannot be filled meaningfully, stop and return:
           {{"error": "Content too short for the template. Please provide more detailed content."}}
        
        ### Output Requirements
        - Return **only a valid JSON object**.
        - <replacementContent> key contains the mapping of placeholders to rewritten text.
        - <iconIntents> key (if applicable) contains an array of icon intent strings.
        - Keys = exact placeholder text from the provided list.
        - Values = strings or string arrays depending on placeholder type.
        - No explanations, markdown, or extra commentary.
        - Ensure the JSON is syntactically valid.
        """
    else:
        prompt = f"""You are a highly precise PowerPoint content mapper. 
        Your task is to map the provided content directly to the given placeholders **exactly as written**, without rewriting or rephrasing it — 
        but with logical splitting and proper assignment based on context.
        
        If you cannot produce a valid mapping for every placeholder, 
        return only this JSON:
        {{"error": "Content too short for the template. Please provide more detailed content."}}
        
        ### Inputs
        - Content: {req.content}
        - Placeholders: {json.dumps(textBoxList, indent=2)}
        
        ### Mapping Rules
        1. **Do not rewrite or rephrase** the text; only split or assign it logically.
        2. Identify the intent of each placeholder (e.g., “Discover”, “Plan”, “Create”, “Deliver”).
        3. Group related sentences, phrases, or bullet points from the content and assign them to the most contextually relevant placeholder.
           - Example: All lines mentioning “research”, “analysis”, or “identifying needs” → map to “Discover”.
           - Lines about “strategy”, “planning”, “goal setting” → map to “Plan”.
           - Lines about “design”, “development”, “execution” → map to “Create”.
           - Lines about “testing”, “delivery”, “measurement”, “results” → map to “Deliver”.
        4. If placeholder type is "list", return an array of bullet points (minimum 1 and maximum as per the template image provided). When creating bullet lists, use clean text without adding any extra symbols such as hyphens (-), asterisks (*), or other bullet markers — return plain text items only. If placeholder type is "text", return a single concise string.
        5. **Never duplicate the same sentences** across placeholders unless the content itself repeats them exactly.
        6. Maintain the factual meaning and original order of ideas wherever possible.
        7. If any placeholder cannot be filled with meaningful data, stop and return:
           {{"error": "Content too short for the template. Please provide more detailed content."}}
        8. If the provided content clearly and meaningfully fills only a subset of placeholders (for example, a title and four main sections), this is acceptable. Do not force-fill empty placeholders with guesses or duplicated text. Only leave placeholders empty if no relevant content exists for them.
        9. if iconCount->({len(iconsCount)}) is greater than 0 then for each section you generate content for also generate icon intents for each icon on the slide as per the following:
        10. There are {len(iconsCount)} icons on the slide; do not reference or describe them in the text. in the json at the end just give a key called iconIntents with a list of {len(iconsCount)} strings describing the intended meaning or concept for each icon in order for each section determining it by the heading and content you generated.
        11. if iconCount->({len(iconsCount)}) is 0 then do not add iconIntents key in the final json.
        12. if iconCount->({len(iconsCount)}) is greater than 0 then give an array of {len(iconsCount)} strings for iconIntents in a json key <iconIntents> and the rest of the content in json key <replacementContent>.
        13. if iconCount->({len(iconsCount)}) is 0 then give only the content in json key <replacementContent>.
        ### Output Format
        - Output **strictly valid JSON only**.
        - <replacementContent> key contains the mapping of placeholders to rewritten text.
        - <iconIntents> key (if applicable) contains an array of icon intent strings.
        - Keys = exact placeholder text from the provided list.
        - Values = strings or string arrays depending on placeholder type.
        - No markdown, no explanations, no extra commentary.
        """

    image_base64 = download_image(req.imageUrl)
    response = client.chat.completions.create(
        model="gpt-5-mini",
        response_format={"type": "json_object"},
        messages=[
            {
                "role": "system",
                "content": "You are a precise PowerPoint content generator that strictly outputs valid JSON."
            },
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": image_base64
                        }
                    }
                ]
            }
        ],
        temperature=0.7
    )
    raw_output = response.choices[0].message.content
    cleanedJson = json.loads(raw_output.strip("`").replace("json","",1).strip())
    
    print("\n----- Prompted ----",prompt,"\n---end prompt---","\n------\nGenerated JSON:", cleanedJson,"\n------\n")
    if "iconIntents" in cleanedJson:
        print("Icons to be replaced:", cleanedJson["iconIntents"])
        iconIntents = cleanedJson.pop("iconIntents")
        for i in range(len(iconsCount)):
            print("Replacing icon with intent:", iconIntents[i])
            replace_icons(pptx_reference, 0, iconsCount[i]['shape'], get_icon(iconIntents[i])['result']['hits'][0]['_id'])
    if "replacementContent" in cleanedJson:
        cleanedJson = cleanedJson["replacementContent"]

    if not validateJson(cleanedJson, textBoxList):
        if os.path.exists(pptx_path):
            os.remove(pptx_path)
        return {"error": "Error Generating PPTX, Content too short for the template. Please provide more detailed content."}
    else:
        updated_pptx = updateTemplatePlaceholders(pptx_reference, 0, cleanedJson)

        # Step 3: Generate unique filename
        unique_id = uuid.uuid4().hex[:8]  # short UUID
        public_filename = f"presentation_{unique_id}.pptx"
        public_path = os.path.join(GENERATED_DIR, public_filename)

        # Move to public folder
        shutil.copy(updated_pptx, public_path)
        # make the file publicly readable
        os.chmod(public_path, 0o755)
        # Delete the temporary file
        if os.path.exists(updated_pptx):
            os.remove(updated_pptx)

        if os.path.exists(pptx_path):
            os.remove(pptx_path)

        # Step 4: Return public URL
        file_url = f"{DOMAIN_NAME}api/{GENERATED_DIR}/{public_filename}"
        return {"file_url": file_url}

@app.post("/upload-files/")
async def upload_files(files: List[UploadFile] = File(...)):
    saved_files = []

    for file in files:
        # Always use the original filename
        filename = file.filename
        file_path = os.path.join(UPLOAD_DIR, filename)

        # "wb" mode automatically replaces file if it already exists
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Make the file publicly readable
        os.chmod(file_path, 0o755)

        # Build file URL
        file_url = f"{DOMAIN_NAME}api/{UPLOAD_DIR}/{filename}"
        saved_files.append({"filename": filename, "url": file_url})

    return {"uploaded": saved_files}

@app.post("/consolidate-slides")
def consolidate_slides(payload: ConsolidateSlidesRequest):
    if not payload.ppt_paths:
        raise HTTPException(status_code=400, detail="ppt_paths cannot be empty")

    final_prs = None
    i = 0
    for ppt_path in payload.ppt_paths:
        x = ""
        if(os.path.exists(os.path.join(GENERATED_DIR, *ppt_path.split('/')[-1:]))):
            x = os.path.join(GENERATED_DIR, *ppt_path.split('/')[-1:])
        elif(os.path.exists(os.path.join(UPLOAD_DIR, *ppt_path.split('/')[-1:]))):
            x = os.path.join(UPLOAD_DIR, *ppt_path.split('/')[-1:])
        else:
            raise HTTPException(
                status_code=404,
                detail=f"PPT not found: {x}"
            )

        src_prs = Presentation(x)

        if final_prs is None:
            final_prs = Presentation(x)
            # remove_all_slides(final_prs)
            final_prs.slide_width = src_prs.slide_width
            final_prs.slide_height = src_prs.slide_height
            # remove_all_slides(final_prs)

        if i != 0:
            for slide in src_prs.slides:
                copy_slide(slide, final_prs)
        i += 1

    output_filename = f"consolidated_{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(GENERATED_DIR, output_filename)

    final_prs.save(output_path)
    file_url = f"{DOMAIN_NAME}api/{GENERATED_DIR}/{output_filename}"
    return {
        "message": "Slides consolidated successfully",
        "output_path": file_url,
        "total_slides": len(final_prs.slides)
    }

def remove_all_slides(prs):
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)

